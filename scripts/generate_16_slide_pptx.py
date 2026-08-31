"""
scripts/generate_16_slide_pptx.py
Generates an attractive, professional 16-slide PowerPoint Presentation (.pptx)
for the IDS Forge COM4901 Final Year Research Defence / Viva Examination.
Output: docs/Viva_Presentation_IDS_Forge_14519.pptx
"""

import os
import sys
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

def create_presentation(output_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette Constants
    NAVY_DARK = RGBColor(15, 23, 42)       # #0F172A - Title / Headers
    BLUE_PRIMARY = RGBColor(37, 99, 235)   # #2563EB - Accent / Highlights
    BLUE_DARK = RGBColor(30, 58, 138)      # #1E3A8A - Subtitles
    SLATE_TEXT = RGBColor(51, 65, 85)      # #334155 - Body Text
    BG_CARD = RGBColor(248, 250, 252)      # #F8FAFC - Card Background
    BORDER_CARD = RGBColor(226, 232, 240)  # #E2E8F0 - Card Border
    EMERALD = RGBColor(16, 185, 129)       # #10B981 - Success Green
    WHITE = RGBColor(255, 255, 255)

    def add_header(slide, title_text, category_text="KIU COM4901 FINAL YEAR DEFENCE"):
        # Category Tag
        tx_box_cat = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
        tf_cat = tx_box_cat.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.name = 'Calibri'
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = BLUE_PRIMARY

        # Title Text
        tx_box_t = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
        tf_t = tx_box_t.text_frame
        tf_t.word_wrap = True
        p_t = tf_t.paragraphs[0]
        p_t.text = title_text
        p_t.font.name = 'Calibri'
        p_t.font.size = Pt(22)
        p_t.font.bold = True
        p_t.font.color.rgb = NAVY_DARK

    def add_card(slide, left, top, width, height, bg_color=BG_CARD, border_color=BORDER_CARD):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
        return shape

    def add_notes(slide, notes_text):
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        tf.text = notes_text

    blank_layout = prs.slide_layouts[6]

    # =========================================================================
    # SLIDE 1: TITLE SLIDE (COVER)
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_layout)
    bg1 = add_card(slide1, Inches(0.5), Inches(0.5), Inches(12.333), Inches(6.5), NAVY_DARK, NAVY_DARK)

    tx_t1 = slide1.shapes.add_textbox(Inches(1.0), Inches(1.2), Inches(11.333), Inches(1.8))
    tf1 = tx_t1.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    p1.text = "IDS FORGE ⚒️"
    p1.font.name = 'Calibri'
    p1.font.size = Pt(40)
    p1.font.bold = True
    p1.font.color.rgb = EMERALD

    p2 = tf1.add_paragraph()
    p2.text = "A Machine Learning-Based Hybrid Intrusion Detection System for IoT Networks"
    p2.font.name = 'Calibri'
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = WHITE

    card_meta = add_card(slide1, Inches(1.0), Inches(3.4), Inches(11.333), Inches(3.0), RGBColor(30, 41, 59), RGBColor(51, 65, 85))
    tx_m = slide1.shapes.add_textbox(Inches(1.3), Inches(3.6), Inches(10.7), Inches(2.6))
    tf_m = tx_m.text_frame
    tf_m.word_wrap = True

    meta_items = [
        ("Student Name:", "R.M.L.S.B. Wijerathna (Student ID: 14519)"),
        ("Degree Program:", "BSc (Hons) in Computer Networks and Cyber Security"),
        ("Module Title:", "COM4901 - Final Year Individual Project"),
        ("Project Supervisor:", "Mr. Sahan Weerasinghe"),
        ("Institution:", "KIU University, Sri Lanka"),
        ("GitHub Repository:", "https://github.com/lakmal6214/IDS-Forge-Research")
    ]
    for lbl, val in meta_items:
        pm = tf_m.add_paragraph()
        r_lbl = pm.add_run()
        r_lbl.text = f"{lbl:<22} "
        r_lbl.font.bold = True
        r_lbl.font.size = Pt(13)
        r_lbl.font.color.rgb = BLUE_PRIMARY
        
        r_val = pm.add_run()
        r_val.text = val
        r_val.font.size = Pt(13)
        r_val.font.color.rgb = WHITE

    add_notes(slide1, "Good morning respected supervisor Mr. Sahan Weerasinghe and panel members. Welcome to my final defense presentation for COM4901. My project is IDS Forge: A Machine Learning-Based Hybrid Intrusion Detection System for Internet of Things Networks.")

    # =========================================================================
    # SLIDE 2: EXECUTIVE SUMMARY & HIGHLIGHTS
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_layout)
    add_header(slide2, "Executive Summary & Core Benchmark Highlights")

    # 4 Metric Cards
    metrics = [
        ("100.00%", "Classification Accuracy", "Zero false positives across 3,000 BoT-IoT test flows", BLUE_PRIMARY),
        ("0.034 ms", "Per-Packet Latency", "Microsecond execution speed ideal for IoT edge gateways", EMERALD),
        ("214.7 MB", "RAM Footprint", "Consumes <6.2% CPU & 214 MB RAM on edge hardware", BLUE_DARK),
        ("100.00%", "Zero-Day Fallback Recall", "Intercepts novel attack probes bypassed by rules", NAVY_DARK)
    ]
    col_w = Inches(2.8)
    gap = Inches(0.3)
    start_x = Inches(0.8)

    for idx, (val, title, desc, color) in enumerate(metrics):
        cx = start_x + idx * (col_w + gap)
        add_card(slide2, cx, Inches(1.6), col_w, Inches(5.2), BG_CARD, BORDER_CARD)
        
        tx = slide2.shapes.add_textbox(cx + Inches(0.2), Inches(1.8), col_w - Inches(0.4), Inches(4.8))
        tf = tx.text_frame
        tf.word_wrap = True
        
        p_v = tf.paragraphs[0]
        p_v.text = val
        p_v.font.name = 'Calibri'
        p_v.font.size = Pt(32)
        p_v.font.bold = True
        p_v.font.color.rgb = color
        p_v.alignment = PP_ALIGN.CENTER
        
        p_t = tf.add_paragraph()
        p_t.text = title
        p_t.font.name = 'Calibri'
        p_t.font.size = Pt(14)
        p_t.font.bold = True
        p_t.font.color.rgb = NAVY_DARK
        p_t.alignment = PP_ALIGN.CENTER
        
        tf.add_paragraph().text = "" # spacing
        
        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.name = 'Calibri'
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = SLATE_TEXT
        p_d.alignment = PP_ALIGN.CENTER

    add_notes(slide2, "To summarize my project achievements upfront: IDS Forge achieves 100% accuracy and 100% zero-day threat recall while processing packets in 0.034 milliseconds and consuming only 214 megabytes of RAM.")

    # =========================================================================
    # SLIDE 3: BACKGROUND & MOTIVATION (THE IOT THREAT LANDSCAPE)
    # =========================================================================
    slide3 = prs.slides.add_slide(blank_layout)
    add_header(slide3, "Background & Motivation: The Growing IoT Threat Landscape")

    c1 = add_card(slide3, Inches(0.8), Inches(1.6), Inches(5.7), Inches(5.2))
    tx1 = slide3.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.3), Inches(4.8))
    tf1 = tx1.text_frame
    tf1.word_wrap = True
    
    p = tf1.paragraphs[0]
    p.text = "1. Exponential IoT Ecosystem Growth"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = BLUE_PRIMARY
    
    bullets1 = [
        "Global IoT endpoints expected to exceed 30 Billion devices by 2030.",
        "Deployed across smart cities, industrial SCADA, healthcare telemetry, and home automation.",
        "Generates massive distributed network telemetry across public and private IP networks."
    ]
    for b in bullets1:
        p_b = tf1.add_paragraph()
        p_b.text = f"• {b}"
        p_b.font.size = Pt(12)
        p_b.font.color.rgb = SLATE_TEXT

    p_m = tf1.add_paragraph()
    p_m.text = "\n2. Severe Hardware Constraints"
    p_m.font.bold = True
    p_m.font.size = Pt(16)
    p_m.font.color.rgb = BLUE_PRIMARY

    bullets1_2 = [
        "Embedded microcontrollers (ARM Cortex-M, ESP32, MIPS32) run at modest 80–400 MHz clock rates.",
        "Extremely low volatile RAM memory (<512 MB) and battery power limitations."
    ]
    for b in bullets1_2:
        p_b = tf1.add_paragraph()
        p_b.text = f"• {b}"
        p_b.font.size = Pt(12)
        p_b.font.color.rgb = SLATE_TEXT

    c2 = add_card(slide3, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.2))
    tx2 = slide3.shapes.add_textbox(Inches(7.0), Inches(1.8), Inches(5.3), Inches(4.8))
    tf2 = tx2.text_frame
    tf2.word_wrap = True

    p = tf2.paragraphs[0]
    p.text = "3. Native Vulnerabilities & Botnets"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = BLUE_PRIMARY

    bullets2 = [
        "Plain-text protocols: Unencrypted MQTT (1883), CoAP (5683), raw TCP/UDP.",
        "Hardcoded default credentials (Telnet 23, SSH 22) and unpatched firmware.",
        "Malware Botnets (Mirai, Bashlite, Hajime) compromise IoT endpoints into botnet armies to launch terabit-scale DDoS SYN/UDP floods (>1.2 Tbps)."
    ]
    for b in bullets2:
        p_b = tf2.add_paragraph()
        p_b.text = f"• {b}"
        p_b.font.size = Pt(12)
        p_b.font.color.rgb = SLATE_TEXT

    p_m = tf2.add_paragraph()
    p_m.text = "\n4. Infeasibility of Legacy NIDS"
    p_m.font.bold = True
    p_m.font.size = Pt(16)
    p_m.font.color.rgb = BLUE_PRIMARY

    bullets2_2 = [
        "Legacy enterprise NIDS (Snort, Suricata) rely on heavy Deep Packet Inspection (DPI) & tens of thousands of signatures.",
        "Causes high CPU saturation, out-of-memory kernel panics, and packet buffer overflow on edge gateways."
    ]
    for b in bullets2_2:
        p_b = tf2.add_paragraph()
        p_b.text = f"• {b}"
        p_b.font.size = Pt(12)
        p_b.font.color.rgb = SLATE_TEXT

    add_notes(slide3, "The Internet of Things is expanding rapidly, but IoT devices possess weak microcontrollers and restricted RAM. Malware like Mirai exploits default passwords to build massive botnets, while enterprise security like Snort is too heavy for IoT hardware.")

    # =========================================================================
    # SLIDE 4: PROBLEM STATEMENT (WHY EXISTING SOLUTIONS FAIL)
    # =========================================================================
    slide4 = prs.slides.add_slide(blank_layout)
    add_header(slide4, "Problem Statement: Limitations of Existing Approaches")

    probs = [
        ("1. Pure Signature-Based IDS (S-IDS)", "Fast string matching for known rules, but 0.00% detection recall against zero-day novel attacks & mutated malware variants.", RGBColor(220, 38, 38)),
        ("2. Pure Anomaly-Based IDS (A-IDS)", "Detects novel zero-day threats, but incurs heavy per-packet matrix calculation, high CPU overhead, and high False Positive Rates (FPR).", RGBColor(217, 119, 6)),
        ("3. Inefficient Parallel Hybrid Models", "Existing hybrid systems run signature matching and ML inference in parallel, processing every packet twice and wasting edge gateway CPU cycles.", NAVY_DARK)
    ]

    for idx, (title, desc, color) in enumerate(probs):
        top_y = Inches(1.6 + idx * 1.8)
        add_card(slide4, Inches(0.8), top_y, Inches(11.7), Inches(1.5))
        
        tx = slide4.shapes.add_textbox(Inches(1.1), top_y + Inches(0.15), Inches(11.1), Inches(1.2))
        tf = tx.text_frame
        tf.word_wrap = True
        
        pt = tf.paragraphs[0]
        pt.text = title
        pt.font.bold = True
        pt.font.size = Pt(16)
        pt.font.color.rgb = color
        
        pd = tf.add_paragraph()
        pd.text = desc
        pd.font.size = Pt(13)
        pd.font.color.rgb = SLATE_TEXT

    add_notes(slide4, "This diagram summarizes the problem: Signature engines fail on zero-day attacks, pure Machine Learning engines overload edge CPUs, and parallel hybrid models waste processing power by running both models on every packet.")

    # =========================================================================
    # SLIDE 5: RESEARCH AIM & CORE OBJECTIVES
    # =========================================================================
    slide5 = prs.slides.add_slide(blank_layout)
    add_header(slide5, "Research Aim & Specific Objectives")

    # Aim Banner
    c_aim = add_card(slide5, Inches(0.8), Inches(1.6), Inches(11.7), Inches(1.2), RGBColor(238, 242, 255), BLUE_PRIMARY)
    tx_aim = slide5.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(11.3), Inches(1.0))
    tf_aim = tx_aim.text_frame
    tf_aim.word_wrap = True
    pa = tf_aim.paragraphs[0]
    pa.text = "PRIMARY RESEARCH AIM:"
    pa.font.bold = True
    pa.font.size = Pt(12)
    pa.font.color.rgb = BLUE_PRIMARY
    
    pa2 = tf_aim.add_paragraph()
    pa2.text = "To design, implement, benchmark, and evaluate IDS Forge: a lightweight Machine Learning-Based Hybrid Intrusion Detection System tailored for resource-constrained IoT networks that maximizes accuracy and zero-day threat defense while preserving sub-millisecond latency and minimal memory overhead."
    pa2.font.bold = True
    pa2.font.size = Pt(13)
    pa2.font.color.rgb = NAVY_DARK

    # Objectives Grid (2 Columns)
    objs_left = [
        "Obj 1: Conduct literature review on IoT threat vectors & NIDS bottlenecks.",
        "Obj 2: Preprocess & normalize BoT-IoT benchmark traffic dataset (70/30 split).",
        "Obj 3: Execute 3-stage feature selection (reduce 12 features to 8 optimal attributes).",
        "Obj 4: Construct Phase 1 Signature Engine with 9 protocol-specific IoT security rules."
    ]
    objs_right = [
        "Obj 5: Train & benchmark Phase 2 candidate ML models (DT, RF, Gradient Boosting).",
        "Obj 6: Integrate Phase 1 & 2 into a 2-Tier Sequential Hybrid Engine.",
        "Obj 7: Conduct empirical hardware overhead tracking (CPU %, RAM MB, Latency ms).",
        "Obj 8: Validate zero-day attack recall & publish interactive web dashboard."
    ]

    c_ol = add_card(slide5, Inches(0.8), Inches(3.0), Inches(5.7), Inches(3.8))
    tx_ol = slide5.shapes.add_textbox(Inches(1.0), Inches(3.2), Inches(5.3), Inches(3.4))
    tf_ol = tx_ol.text_frame
    tf_ol.word_wrap = True
    for o in objs_left:
        po = tf_ol.add_paragraph()
        po.text = f"✔ {o}"
        po.font.size = Pt(12)
        po.font.color.rgb = SLATE_TEXT
        tf_ol.add_paragraph().text = ""

    c_or = add_card(slide5, Inches(6.8), Inches(3.0), Inches(5.7), Inches(3.8))
    tx_or = slide5.shapes.add_textbox(Inches(7.0), Inches(3.2), Inches(5.3), Inches(3.4))
    tf_or = tx_or.text_frame
    tf_or.word_wrap = True
    for o in objs_right:
        po = tf_or.add_paragraph()
        po.text = f"✔ {o}"
        po.font.size = Pt(12)
        po.font.color.rgb = SLATE_TEXT
        tf_or.add_paragraph().text = ""

    add_notes(slide5, "My primary aim was to build a lightweight hybrid system for IoT edge networks. I formulated eight specific objectives covering feature selection, rule engineering, machine learning training, sequential integration, and hardware resource evaluation.")

    # =========================================================================
    # SLIDE 6: PROPOSED SYSTEM ARCHITECTURE (TWO-TIER SEQUENTIAL PIPELINE)
    # =========================================================================
    slide6 = prs.slides.add_slide(blank_layout)
    add_header(slide6, "Proposed System Architecture: Two-Tier Sequential Dataflow")

    flow_cards = [
        ("Step 1: Network Ingress", "Raw IoT Traffic Packets arrive at gateway network interface.", BLUE_PRIMARY),
        ("Step 2: Preprocessing", "Feature extraction & Min-Max normalization [0, 1].", BLUE_DARK),
        ("Step 3: Phase 1 Signature", "Evaluates 9 deterministic boolean rules. If matched -> IMMEDIATE BLOCK.", EMERALD),
        ("Step 4: Phase 2 ML Anomaly", "Unmatched flows evaluated by Random Forest (8 features) -> THREAT DECISION.", NAVY_DARK)
    ]

    for idx, (title, desc, color) in enumerate(flow_cards):
        cx = Inches(0.8 + idx * 2.95)
        add_card(slide6, cx, Inches(1.6), Inches(2.7), Inches(5.2))
        
        tx = slide6.shapes.add_textbox(cx + Inches(0.15), Inches(1.8), Inches(2.4), Inches(4.8))
        tf = tx.text_frame
        tf.word_wrap = True
        
        pt = tf.paragraphs[0]
        pt.text = title
        pt.font.bold = True
        pt.font.size = Pt(15)
        pt.font.color.rgb = color
        pt.alignment = PP_ALIGN.CENTER
        
        tf.add_paragraph().text = ""
        
        pd = tf.add_paragraph()
        pd.text = desc
        pd.font.size = Pt(12)
        pd.font.color.rgb = SLATE_TEXT
        pd.alignment = PP_ALIGN.CENTER

    add_notes(slide6, "Here is the architectural core of IDS Forge: Sequential Execution. Known attack packets match Phase 1 rules and are blocked instantly in sub-milliseconds, bypassing ML inference. Only unmatched ambiguous packets pass to Phase 2 Machine Learning.")

    # =========================================================================
    # SLIDE 7: 3-STAGE FEATURE SELECTION PIPELINE
    # =========================================================================
    slide7 = prs.slides.add_slide(blank_layout)
    add_header(slide7, "3-Stage Feature Selection Pipeline")

    stages = [
        ("Stage 1: Pearson Correlation Filter", "Linear Relationship Screening", "Evaluates correlation coefficient r between feature X_i and target label Y: r = \\frac{\\sum (X_i - \\bar{X}_i)(Y - \\bar{Y})}{\\sqrt{\\sum (X_i - \\bar{X}_i)^2 \\sum (Y - \\bar{Y})^2}}. Discards features with r < 0.40.", BLUE_PRIMARY),
        ("Stage 2: Information Gain / Mutual Info", "Non-Linear Dependency Screening", "Quantifies information gain I(X; Y) = \\sum p(x,y) \\log \\frac{p(x,y)}{p(x)p(y)}. Measures entropy reduction H(Y) - H(Y|X) when splitting data on attribute X.", BLUE_DARK),
        ("Stage 3: Recursive Feature Elimination (RFE)", "Wrapper Model Selection", "Fits Random Forest estimator recursively, pruning least important features until top K=8 attributes remain. Reduces feature dimensionality by 33.3%.", EMERALD)
    ]

    for idx, (stitle, ssub, sdesc, color) in enumerate(stages):
        top_y = Inches(1.6 + idx * 1.8)
        add_card(slide7, Inches(0.8), top_y, Inches(11.7), Inches(1.5))
        
        tx = slide7.shapes.add_textbox(Inches(1.1), top_y + Inches(0.12), Inches(11.1), Inches(1.2))
        tf = tx.text_frame
        tf.word_wrap = True
        
        pt = tf.paragraphs[0]
        pt.text = f"{stitle} — {ssub}"
        pt.font.bold = True
        pt.font.size = Pt(15)
        pt.font.color.rgb = color
        
        pd = tf.add_paragraph()
        pd.text = sdesc
        pd.font.size = Pt(11.5)
        pd.font.color.rgb = SLATE_TEXT

    add_notes(slide7, "To keep the model lightweight for edge hardware, I developed a 3-Stage Feature Selection Pipeline combining Pearson linear correlation, Mutual Information non-linear gain, and Recursive Feature Elimination to pick the top 8 optimal features.")

    # =========================================================================
    # SLIDE 8: SELECTED 8 FEATURE ATTRIBUTES DICTIONARY
    # =========================================================================
    slide8 = prs.slides.add_slide(blank_layout)
    add_header(slide8, "Selected 8 Optimal Feature Attributes Dictionary")

    c_tbl = add_card(slide8, Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.2))
    tx_t = slide8.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.3), Inches(4.8))
    tf_t = tx_t.text_frame
    tf_t.word_wrap = True

    p = tf_t.paragraphs[0]
    p.text = "Top 8 Selected Features (Reduced from 12 Attributes via 3-Stage Selection):"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = NAVY_DARK

    feats = [
        ("1. N_IN_Conn_P_DstIP", "Number of inbound connections per destination IP window (r = 0.8265, MI = 0.6729)"),
        ("2. N_IN_Conn_P_SrcIP", "Number of inbound connections per source IP window (r = 0.8066, MI = 0.6688)"),
        ("3. max", "Maximum frame/packet duration recorded within flow window (r = 0.7847, MI = 0.5835)"),
        ("4. srate", "Source packet transmission rate in packets per second (r = 0.6093, MI = 0.5514)"),
        ("5. mean", "Arithmetic mean of flow packet durations (r = 0.7332, MI = 0.5410)"),
        ("6. stddev", "Standard deviation of packet inter-arrival times (r = 0.7367, MI = 0.5358)"),
        ("7. state_number", "Integer encoding connection state SYN, ESTABLISHED, FIN (r = 0.4077, MI = 0.4488)"),
        ("8. dport", "Destination port number identifying target network service (r = 0.6126, MI = 0.3768)")
    ]

    for f_name, f_desc in feats:
        pf = tf_t.add_paragraph()
        r1 = pf.add_run()
        r1.text = f"{f_name:<24} "
        r1.font.bold = True
        r1.font.size = Pt(11.5)
        r1.font.color.rgb = BLUE_PRIMARY
        
        r2 = pf.add_run()
        r2.text = f"— {f_desc}"
        r2.font.size = Pt(11.5)
        r2.font.color.rgb = SLATE_TEXT

    add_notes(slide8, "These 8 features represent connection density, flow duration statistics, sending rate, and destination port. Discarding sport, drate, proto, and min reduced feature space by 33.3% without losing classification accuracy.")

    # =========================================================================
    # SLIDE 9: PHASE 1 DETERMINISTIC SIGNATURE ENGINE
    # =========================================================================
    slide9 = prs.slides.add_slide(blank_layout)
    add_header(slide9, "Phase 1: Deterministic Signature Engine (9 Protocol Rules)")

    c_sig = add_card(slide9, Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.2))
    tx_s = slide9.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.3), Inches(4.8))
    tf_s = tx_s.text_frame
    tf_s.word_wrap = True

    p = tf_s.paragraphs[0]
    p.text = "Deterministic Boolean Matching Logic & Target Attack Vectors:"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = NAVY_DARK

    s_rules = [
        ("Rule 1 (TCP DDoS Flood):", "proto == 0 AND N_IN_Conn_P_SrcIP >= 50", "DDoS Attack"),
        ("Rule 2 (UDP DDoS Flood):", "proto == 1 AND N_IN_Conn_P_DstIP >= 50", "DDoS Attack"),
        ("Rule 3 (HTTP DoS Flood):", "dport in [80, 8080, 443] AND srate >= 100.0", "DoS Attack"),
        ("Rule 4 (UDP DoS Flood):", "proto == 1 AND drate >= 100.0", "DoS Attack"),
        ("Rule 5 (Mirai Telnet Scan):", "dport == 23", "Reconnaissance / Mirai Botnet"),
        ("Rule 6 (SSH Brute-Force):", "dport == 22", "Reconnaissance / Brute-Force"),
        ("Rule 7 (Port Scan Sweep):", "stddev >= 0.5 AND mean <= 0.5", "Reconnaissance Sweep"),
        ("Rule 8 (FTP Data Theft):", "dport == 21", "Data Exfiltration"),
        ("Rule 9 (Connection Anomaly):", "N_IN_Conn_P_SrcIP > 40 OR N_IN_Conn_P_DstIP > 40", "Generic DoS / Anomaly")
    ]

    for r_id, r_cond, r_tgt in s_rules:
        pr = tf_s.add_paragraph()
        r1 = pr.add_run()
        r1.text = f"{r_id:<28} "
        r1.font.bold = True
        r1.font.size = Pt(11)
        r1.font.color.rgb = BLUE_PRIMARY

        r2 = pr.add_run()
        r2.text = f"Condition: {r_cond:<48} "
        r2.font.size = Pt(11)
        r2.font.color.rgb = SLATE_TEXT

        r3 = pr.add_run()
        r3.text = f"[{r_tgt}]"
        r3.font.bold = True
        r3.font.size = Pt(11)
        r3.font.color.rgb = EMERALD

    add_notes(slide9, "Phase 1 implements 9 protocol-specific boolean rules. String matching runs in O(N) linear time, blocking high-density TCP/UDP floods, Mirai Telnet scans, and SSH brute-force probes instantly.")

    # =========================================================================
    # SLIDE 10: PHASE 2 MACHINE LEARNING ANOMALY DETECTOR
    # =========================================================================
    slide10 = prs.slides.add_slide(blank_layout)
    add_header(slide10, "Phase 2: Machine Learning Anomaly Detector Selection")

    models = [
        ("Decision Tree (DT)", "Fast microsecond inference (0.000495 ms), but suffers from structural overfitting on fluctuating live network streams.", RGBColor(217, 119, 6)),
        ("Random Forest (RF - 50 Trees)", "SELECTED MODEL: Ensemble bootstrap aggregation (bagging) averages 50 de-correlated trees, suppressing noise and delivering 100% accuracy with 0.010 ms latency.", EMERALD),
        ("Gradient Boosting (GBM)", "Sequential boosting minimizes residual loss, achieving high accuracy but requiring 4x longer training time (0.424 s) & higher inference overhead.", NAVY_DARK)
    ]

    for idx, (mtitle, mdesc, color) in enumerate(models):
        top_y = Inches(1.6 + idx * 1.8)
        add_card(slide10, Inches(0.8), top_y, Inches(11.7), Inches(1.5))
        
        tx = slide10.shapes.add_textbox(Inches(1.1), top_y + Inches(0.15), Inches(11.1), Inches(1.2))
        tf = tx.text_frame
        tf.word_wrap = True
        
        pt = tf.paragraphs[0]
        pt.text = mtitle
        pt.font.bold = True
        pt.font.size = Pt(16)
        pt.font.color.rgb = color
        
        pd = tf.add_paragraph()
        pd.text = mdesc
        pd.font.size = Pt(12)
        pd.font.color.rgb = SLATE_TEXT

    add_notes(slide10, "For Phase 2 anomaly detection, I evaluated Decision Tree, Random Forest, and Gradient Boosting. Random Forest with 50 decision trees was selected because ensemble bagging eliminates decision noise and prevents overfitting.")

    # =========================================================================
    # SLIDE 11: EXPERIMENTAL SETUP & TECHNOLOGY STACK
    # =========================================================================
    slide11 = prs.slides.add_slide(blank_layout)
    add_header(slide11, "Experimental Setup & Technology Stack")

    c1 = add_card(slide11, Inches(0.8), Inches(1.6), Inches(5.7), Inches(5.2))
    tx1 = slide11.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.3), Inches(4.8))
    tf1 = tx1.text_frame
    tf1.word_wrap = True
    
    p = tf1.paragraphs[0]
    p.text = "Benchmark Dataset Specifications"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = BLUE_PRIMARY

    d_items = [
        "Dataset Schema: BoT-IoT Benchmark Dataset (Koroniotis et al., 2019).",
        "Total Sample Count: 10,000 network traffic flow records.",
        "Data Split: 70% Training Set (7,000 samples) / 30% Testing Set (3,000 samples).",
        "Target Attack Classes: DoS, DDoS, Reconnaissance, Data Theft / Mirai Botnet."
    ]
    for d in d_items:
        pd = tf1.add_paragraph()
        pd.text = f"• {d}"
        pd.font.size = Pt(12)
        pd.font.color.rgb = SLATE_TEXT

    c2 = add_card(slide11, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.2))
    tx2 = slide11.shapes.add_textbox(Inches(7.0), Inches(1.8), Inches(5.3), Inches(4.8))
    tf2 = tx2.text_frame
    tf2.word_wrap = True

    p = tf2.paragraphs[0]
    p.text = "Software Architecture & Libraries"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = BLUE_PRIMARY

    t_items = [
        "Programming Language: Python 3.14 (Virtual Environment .venv).",
        "Machine Learning: scikit-learn 1.7.1 (DT, RF, GBM, RFE).",
        "Data Manipulation: pandas 2.2.3, numpy 2.2.3.",
        "Hardware Profiling: psutil 7.0.0 (RAM RSS & CPU % tracking).",
        "Interactive Dashboard: Streamlit 1.61.1 web UI framework."
    ]
    for t in t_items:
        pt = tf2.add_paragraph()
        pt.text = f"• {t}"
        pt.font.size = Pt(12)
        pt.font.color.rgb = SLATE_TEXT

    add_notes(slide11, "The experimental evaluation uses the benchmark BoT-IoT dataset schema across 10,000 traffic flows split 70/30. The software stack was built entirely in Python 3.14 using scikit-learn, pandas, psutil, and Streamlit.")

    # =========================================================================
    # SLIDE 12: EMPIRICAL CLASSIFICATION PERFORMANCE RESULTS
    # =========================================================================
    slide12 = prs.slides.add_slide(blank_layout)
    add_header(slide12, "Empirical Classification Performance Benchmark Results")

    c_res = add_card(slide12, Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.2))
    tx_r = slide12.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.3), Inches(4.8))
    tf_r = tx_r.text_frame
    tf_r.word_wrap = True

    p = tf_r.paragraphs[0]
    p.text = "Comparative Classification Performance Across System Configurations:"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = NAVY_DARK

    res_data = [
        ("Standalone Signature IDS", "98.27%", "97.11%", "0.000%", "0.0275 ms", "No (0% Zero-Day Recall)"),
        ("Standalone ML Anomaly IDS", "100.00%", "100.00%", "0.000%", "0.0101 ms", "Yes"),
        ("IDS Forge Hybrid (Proposed)", "100.00%", "100.00%", "0.000%", "0.0340 ms", "Yes (100% Zero-Day Recall)")
    ]

    for sys_name, acc, dr, fpr, lat, zd in res_data:
        pr = tf_r.add_paragraph()
        tf_r.add_paragraph().text = ""
        
        r1 = pr.add_run()
        r1.text = f"• {sys_name:<30} "
        r1.font.bold = True
        r1.font.size = Pt(13)
        r1.font.color.rgb = BLUE_PRIMARY if "Proposed" not in sys_name else EMERALD

        r2 = pr.add_run()
        r2.text = f"Accuracy: {acc}  |  Recall: {dr}  |  FPR: {fpr}  |  Latency: {lat}"
        r2.font.size = Pt(12)
        r2.font.color.rgb = SLATE_TEXT

    add_notes(slide12, "Looking at empirical classification results across 3,000 test flows: Standalone signature rules achieved 98.27% accuracy. Our proposed IDS Forge Hybrid achieved 100% classification accuracy and 100% detection rate with a 0% false positive rate.")

    # =========================================================================
    # SLIDE 13: HARDWARE OVERHEAD & RESOURCE BENCHMARKING
    # =========================================================================
    slide13 = prs.slides.add_slide(blank_layout)
    add_header(slide13, "Hardware Resource Overhead & Edge Feasibility")

    hw_metrics = [
        ("0.034 ms", "Average Per-Packet Latency", "Microsecond execution speed enables real-time line-rate packet inspection.", BLUE_PRIMARY),
        ("6.2%", "Average CPU Load", "Saves >68% CPU cycles compared to running standalone DPI engines.", EMERALD),
        ("214.70 MB", "RAM Memory Footprint", "Fits comfortably within low-cost 512 MB IoT gateway memory limits.", NAVY_DARK)
    ]

    for idx, (val, title, desc, color) in enumerate(hw_metrics):
        cx = Inches(0.8 + idx * 3.95)
        add_card(slide13, cx, Inches(1.6), Inches(3.7), Inches(5.2))
        
        tx = slide13.shapes.add_textbox(cx + Inches(0.2), Inches(1.8), Inches(3.3), Inches(4.8))
        tf = tx.text_frame
        tf.word_wrap = True
        
        pv = tf.paragraphs[0]
        pv.text = val
        pv.font.bold = True
        pv.font.size = Pt(36)
        pv.font.color.rgb = color
        pv.alignment = PP_ALIGN.CENTER
        
        pt = tf.add_paragraph()
        pt.text = title
        pt.font.bold = True
        pt.font.size = Pt(14)
        pt.font.color.rgb = NAVY_DARK
        pt.alignment = PP_ALIGN.CENTER
        
        tf.add_paragraph().text = ""
        
        pd = tf.add_paragraph()
        pd.text = desc
        pd.font.size = Pt(12)
        pd.font.color.rgb = SLATE_TEXT
        pd.alignment = PP_ALIGN.CENTER

    add_notes(slide13, "Resource overhead evaluation confirms edge gateway feasibility: Processing latency is 0.034 ms per packet, CPU load remains under 6.2%, and memory footprint is 214.7 MB RAM.")

    # =========================================================================
    # SLIDE 14: ZERO-DAY ATTACK SIMULATION RESULTS
    # =========================================================================
    slide14 = prs.slides.add_slide(blank_layout)
    add_header(slide14, "Zero-Day Attack Detection Simulation Results")

    c_zd = add_card(slide14, Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.2))
    tx_z = slide14.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.3), Inches(4.8))
    tf_z = tx_z.text_frame
    tf_z.word_wrap = True

    p = tf_z.paragraphs[0]
    p.text = "Simulated Rule-Bypass Scenario (Disabling Phase 1 Reconnaissance Rules):"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = NAVY_DARK

    zd_items = [
        ("Denial of Service (DoS):", "1,050 Samples", "100.00% Signature Recall", "100.00% Hybrid Fallback Recall", "Maintained"),
        ("Distributed DoS (DDoS):", "1,050 Samples", "100.00% Signature Recall", "100.00% Hybrid Fallback Recall", "Maintained"),
        ("Reconnaissance (Zero-Day):", "600 Samples", "85.56% Signature Recall", "100.00% Hybrid Fallback Recall", "+14.44% Improvement"),
        ("Data Theft / Mirai Botnet:", "300 Samples", "100.00% Signature Recall", "100.00% Hybrid Fallback Recall", "Maintained")
    ]

    for atk, samp, sig_r, hyb_r, imp in zd_items:
        pz = tf_z.add_paragraph()
        tf_z.add_paragraph().text = ""
        
        r1 = pz.add_run()
        r1.text = f"• {atk:<26} "
        r1.font.bold = True
        r1.font.size = Pt(12)
        r1.font.color.rgb = BLUE_PRIMARY

        r2 = pz.add_run()
        r2.text = f"Rules Recall: {sig_r}  -->  Hybrid Fallback Recall: {hyb_r}  "
        r2.font.size = Pt(12)
        r2.font.color.rgb = SLATE_TEXT

        r3 = pz.add_run()
        r3.text = f"[{imp}]"
        r3.font.bold = True
        r3.font.size = Pt(12)
        r3.font.color.rgb = EMERALD

    add_notes(slide14, "To test zero-day defense, I disabled Phase 1 rules for reconnaissance scans. Signature-only detection dropped to 85.56%, but Phase 2 Random Forest intercepted 100% of bypassed probes, proving +14.44% zero-day recall improvement.")

    # =========================================================================
    # SLIDE 15: SOFTWARE DELIVERABLES & WEB UI DASHBOARD
    # =========================================================================
    slide15 = prs.slides.add_slide(blank_layout)
    add_header(slide15, "Software Deliverables & Interactive Web Dashboard")

    c1 = add_card(slide15, Inches(0.8), Inches(1.6), Inches(5.7), Inches(5.2))
    tx1 = slide15.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.3), Inches(4.8))
    tf1 = tx1.text_frame
    tf1.word_wrap = True
    
    p = tf1.paragraphs[0]
    p.text = "Interactive Streamlit Web App (`app.py`)"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = BLUE_PRIMARY

    sw_items = [
        "Real-Time Packet Inspector & Rule Evaluator.",
        "Interactive ML Model Benchmark Comparison.",
        "Zero-Day Threat Simulation & Fallback Verification.",
        "Automated PDF Report & Metric Log Export."
    ]
    for sw in sw_items:
        ps = tf1.add_paragraph()
        ps.text = f"• {sw}"
        ps.font.size = Pt(12)
        ps.font.color.rgb = SLATE_TEXT

    c2 = add_card(slide15, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.2))
    tx2 = slide15.shapes.add_textbox(Inches(7.0), Inches(1.8), Inches(5.3), Inches(4.8))
    tf2 = tx2.text_frame
    tf2.word_wrap = True

    p = tf2.paragraphs[0]
    p.text = "1-Click Launchers & GitHub Codebase"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = BLUE_PRIMARY

    gh_items = [
        "Windows 1-Click Batch Script: run_ids_forge.bat",
        "Mac / Linux Shell Script: run_ids_forge.sh",
        "Full Modular Package Architecture in src/",
        "Public Repository: https://github.com/lakmal6214/IDS-Forge-Research"
    ]
    for gh in gh_items:
        pg = tf2.add_paragraph()
        pg.text = f"• {gh}"
        pg.font.size = Pt(12)
        pg.font.color.rgb = SLATE_TEXT

    add_notes(slide15, "IDS Forge includes a complete Streamlit web UI dashboard, automated PDF export tools, 1-click batch scripts for Windows, and is fully published on GitHub.")

    # =========================================================================
    # SLIDE 16: CONCLUSION, FUTURE WORK & Q&A
    # =========================================================================
    slide16 = prs.slides.add_slide(blank_layout)
    bg16 = add_card(slide16, Inches(0.5), Inches(0.5), Inches(12.333), Inches(6.5), NAVY_DARK, NAVY_DARK)

    tx16 = slide16.shapes.add_textbox(Inches(1.0), Inches(1.0), Inches(11.333), Inches(5.5))
    tf16 = tx16.text_frame
    tf16.word_wrap = True

    p = tf16.paragraphs[0]
    p.text = "Conclusion & Future Work"
    p.font.bold = True
    p.font.size = Pt(28)
    p.font.color.rgb = EMERALD

    c_items = [
        "Conclusion: IDS Forge proves that combining deterministic signature matching with Random Forest anomaly classification resolves the speed vs. generalization trade-off for IoT edge gateways.",
        "Theoretical Contribution: Formulated a 3-stage feature selection methodology proving 8 statistical attributes are sufficient for multi-category IoT threat detection.",
        "Practical Contribution: Open-source, lightweight sub-millisecond HIDS ready for embedded gateway deployment (<215 MB RAM).",
        "Future Directions: Physical Raspberry Pi 4 edge hardware testing, ONNX microsecond runtime export, and Deep Variational Autoencoders."
    ]

    for c in c_items:
        pc = tf16.add_paragraph()
        tf16.add_paragraph().text = ""
        r1 = pc.add_run()
        r1.text = f"✔ {c}"
        r1.font.size = Pt(13)
        r1.font.color.rgb = WHITE

    tf16.add_paragraph().text = ""
    p_qa = tf16.add_paragraph()
    p_qa.text = "THANK YOU!  |  OPEN FOR QUESTIONS & DISCUSSION"
    p_qa.font.bold = True
    p_qa.font.size = Pt(20)
    p_qa.font.color.rgb = BLUE_PRIMARY
    p_qa.alignment = PP_ALIGN.CENTER

    add_notes(slide16, "In conclusion, IDS Forge successfully solves the trade-off between detection speed and zero-day threat defense for IoT edge networks. Thank you for your time and attention. I am now open for any questions!")

    # Save Presentation
    prs.save(output_path)
    print(f"[+] 16-Slide Viva Presentation (.pptx) generated successfully at: {output_path}")

if __name__ == '__main__':
    out_pptx = os.path.join("docs", "Viva_Presentation_IDS_Forge_14519.pptx")
    os.makedirs("docs", exist_ok=True)
    create_presentation(out_pptx)
