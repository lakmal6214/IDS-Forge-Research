"""
generate_pptx.py
Converts presentation markdown into a widescreen (16:9) PowerPoint presentation (.pptx)
with dark high-tech styling, clear slide layout, and complete speaker notes for all 35 slides.
"""

import os
import re
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation_pptx():
    prs = Presentation()
    # 16:9 Widescreen Layout
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    presentation_path = r"C:\Users\Lakmal\.gemini\antigravity\brain\72e3f9c9-eae3-40d8-85d2-cb806c66bdee\presentation.md"
    if not os.path.exists(presentation_path):
        print(f"Error: {presentation_path} not found.")
        return

    with open(presentation_path, 'r', encoding='utf-8') as f:
        content = f.read()

    # Split slides by "## SLIDE "
    raw_slides = re.split(r'## SLIDE \d+:', content)
    
    blank_slide_layout = prs.slide_layouts[6] # Blank layout

    print(f"[*] Processing {len(raw_slides)-1} Slides into PPTX...")

    for idx, raw in enumerate(raw_slides[1:], 1):
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Set Dark Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0x08, 0x09, 0x11) # Dark navy

        lines = raw.strip().split('\n')
        title_text = lines[0].strip() if lines else f"Slide {idx}"

        body_lines = []
        speaker_notes = ""
        in_notes = False

        for line in lines[1:]:
            line_str = line.strip()
            if line_str.startswith('> **Speaker Notes:**') or line_str.startswith('> Speaker Notes:'):
                in_notes = True
                speaker_notes += line_str.replace('> **Speaker Notes:**', '').replace('> Speaker Notes:', '').strip() + " "
            elif in_notes:
                if line_str.startswith('>'):
                    speaker_notes += line_str.replace('>', '').strip() + " "
                else:
                    in_notes = False
            else:
                if line_str and not line_str.startswith('---'):
                    body_lines.append(line_str)

        # 1. Slide Title Box
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.7), Inches(1.0))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = 'Outfit'
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x00, 0xF0, 0xFF) # Cyan

        # 2. Slide Content Box
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
        cf = content_box.text_frame
        cf.word_wrap = True

        for b_idx, line in enumerate(body_lines):
            p = cf.add_paragraph() if b_idx > 0 else cf.paragraphs[0]
            clean_text = line.replace('**', '').replace('- ', '').replace('* ', '')
            p.text = clean_text
            p.font.name = 'Inter'
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(0xF0, 0xF4, 0xF8)
            p.space_after = Pt(10)
            if line.startswith('- ') or line.startswith('* '):
                p.level = 0

        # 3. Add Speaker Notes
        if speaker_notes:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = speaker_notes.strip()

    output_path = r"c:\Users\Lakmal\Documents\Research\Final_Presentation_14519.pptx"
    prs.save(output_path)
    print(f"[+] Final Widescreen PowerPoint Presentation created successfully at: {output_path}")

if __name__ == '__main__':
    create_presentation_pptx()
